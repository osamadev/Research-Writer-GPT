from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from langchain.agents import Tool
import streamlit as st
import re
import io

class PresentationCreationTool(Tool):
    def __init__(self, name, description):
        super().__init__(name=name, description=description, func=self.create_presentation)

    def create_presentation(self, conversation_history):
        # Generate slide content with LLM
        from openai import OpenAI
        OpenAI.api_key = st.secrets["OPENAI_API_KEY"]
        client = OpenAI()

        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a PowerPoint expert, designer, content organizer, and slides generator."},
                {"role": "user", "content": f"""Format and generate PowerPoint slides using this context:\n{conversation_history}\n
                 Prepare the slides from the given conversation history and break the content logically according to the topic into slides, 
                 each slide should start with a title in that format "slide 1: Title of slide 1" in a separate line. 
                 The title of each slide should be in the required format and in a new line"""}
            ],
            max_tokens=2000,
            temperature=0.3
        )
        generated_content = response.choices[0].message.content

        # Create a presentation object
        prs = Presentation()

        # Define a professional and elegant design
        bg_color = RGBColor(245, 245, 245)  # Light grey background for elegance
        title_font_color = RGBColor(0, 51, 102)  # Dark blue for titles
        content_font_color = RGBColor(32, 32, 32)  # Almost black for content
        font_family = "Arial"  # Professional font family

        # Updated regex to correctly extract slides
        slides = extract_presentation_content(generated_content)

        for slide in slides:
            # Split title and content
            title = slide.get('title')
            content = slide.get('content')
        
            # Add a blank slide
            slide_layout = prs.slide_layouts[6]  # 6 corresponds to blank slide layout
            slide = prs.slides.add_slide(slide_layout)

            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

            # Add title
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.clear()  # Clear any existing content in the frame
            title_frame.word_wrap = True
            title_p = title_frame.add_paragraph()
            title_p.text = title.strip()
            title_p.font.size = Pt(36)
            title_p.font.color.rgb = title_font_color
            title_p.font.name = font_family

            # Add content as bullet list
            content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            content_frame = content_shape.text_frame
            content_frame.clear()  # Clear any existing content in the frame
            content_frame.word_wrap = True
            for line in content.split('\n'):
                p = content_frame.add_paragraph()
                p.text = line.strip()
                p.level = 0  # Adjust as needed for nested bullets
                p.font.size = Pt(20)
                p.font.color.rgb = content_font_color
                p.font.name = font_family

            # Adjust margins
            content_frame.margin_left = Emu(457200)  # Approx. 0.5 inch
            content_frame.margin_right = Emu(457200)  # Approx. 0.5 inch
            content_frame.margin_top = Emu(228600)  # Approx. 0.25 inch
            content_frame.margin_bottom = Emu(228600)  # Approx. 0.25 inch

        # Save to a BytesIO object
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        return ppt_io, generated_content
    
def extract_presentation_content(generated_content):
    slides = re.split(r'Slide \d+:', generated_content)

    # Process each slide (excluding the first element which is text before the first slide)
    slides_data = []
    for slide in slides[1:]:
        # Split the slide into title and body by taking the first sentence as the title
        title_body_split = slide.strip().split(' ', 1)
        if len(title_body_split) == 2:
            title, body = title_body_split
            slides_data.append({'title': title, 'content': body})    
    return slides_data

# Initialize the tool
presentation_tool = PresentationCreationTool(
    name="Presentation Generation Tool",
    description="Creates and generates PowerPoint presentations from conversation history. It will be used when the user requests to generate a presnetation out of his conversation"
)